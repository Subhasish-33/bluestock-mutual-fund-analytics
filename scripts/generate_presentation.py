"""
Module: generate_presentation.py

Part of the Bluestock Mutual Fund Analytics Capstone.
This script provides functionality for the data pipeline and analytics engine.
"""

import collections 
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
import os

def create_presentation():
    # Create presentation object
    prs = Presentation()
    
    # Define slide layouts
    title_slide_layout = prs.slide_layouts[0]
    bullet_slide_layout = prs.slide_layouts[1]
    
    # --- Slide 1: Title ---
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Bluestock Mutual Fund Analytics Capstone"
    subtitle.text = "End-to-End Data Pipeline & Analytics Dashboard\nFinal Presentation"

    # --- Slide 2: Problem & Objective ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "Problem & Objective"
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.text = "Problem Statement:"
    p = tf.add_paragraph()
    p.text = "Evaluating mutual funds requires analyzing historical NAV, understanding risk profiles, and standardizing data from multiple sources."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Objectives:"
    p = tf.add_paragraph()
    p.text = "Build an automated ETL pipeline for daily NAVs."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Design a reliable, local data warehouse (SQLite)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Calculate and visualize key performance metrics."
    p.level = 1

    # --- Slide 3: Data Sources ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "Data Sources"
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.text = "Primary Sources utilized:"
    p = tf.add_paragraph()
    p.text = "AMFI Live API: Daily streaming NAV data."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Historical CSVs: 5-year NAV history."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Metadata & Transactions: Synthesized investor activity and scheme profiles."
    p.level = 1

    # --- Slide 4: Architecture ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "System Architecture"
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.text = "Data Engineering Pipeline:"
    p = tf.add_paragraph()
    p.text = "Extraction: Python Requests & Pandas."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Transformation: Anomaly detection, localized forward filling."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Loading: Localized SQLite relational database."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Automation: Scheduled python jobs for background fetching."
    p.level = 1

    # --- Slide 5: EDA Highlights 1 ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "EDA Highlights: AUM & Volatility"
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.text = "Key Observations:"
    p = tf.add_paragraph()
    p.text = "AUM Skewness: Top 10% of funds control 60%+ of the industry AUM."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Volatility Clusters: Small Cap and Sectoral funds map consistently to higher daily standard deviations."
    p.level = 1

    # --- Slide 6: EDA Highlights 2 ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "EDA Highlights: Trading Days"
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.text = "Data Quality Insights:"
    p = tf.add_paragraph()
    p.text = "Non-trading Days: Discovered systematic missing NAVs due to market holidays."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Resolution: Transitioned to 252-trading day scaling for annualized metrics."
    p.level = 1

    # --- Slide 7: Performance Metrics ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "Standard Performance Metrics"
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.text = "Baseline Evaluation:"
    p = tf.add_paragraph()
    p.text = "CAGR: Compounded annual growth based on trading days."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Sharpe Ratio: Risk-adjusted return using a 6% risk-free rate."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Sortino Ratio: Downside-focused risk adjustment."
    p.level = 1

    # --- Slide 8: Advanced Analytics ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "Advanced Risk & Cohort Metrics"
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.text = "Deep-dive Analytics:"
    p = tf.add_paragraph()
    p.text = "Value at Risk (VaR) & Conditional VaR (CVaR)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Maximum Drawdown analysis."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Sector Concentration (HHI) for diversification assessment."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "SIP Continuity Cohorts to monitor investor retention."
    p.level = 1

    # --- Slide 9: Dashboard Overview ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "Interactive Streamlit Dashboard"
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.text = "Features:"
    p = tf.add_paragraph()
    p.text = "Multi-page navigation (Overview, Funds, Performance, Risk)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Dynamic sidebar filters (Category, Risk, Ratings)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "[Insert Screenshot Here]"
    p.level = 1

    # --- Slide 10: Performance Deep Dive ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "Dashboard Visualizations"
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.text = "Visuals built with Plotly:"
    p = tf.add_paragraph()
    p.text = "Interactive scatter plots mapping Return vs Volatility."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Rolling Sharpe Ratio line charts."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "[Insert Screenshot Here]"
    p.level = 1

    # --- Slide 11: Key Findings & Recommendations ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "Key Findings & Recommendations"
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.text = "Strategic Takeaways:"
    p = tf.add_paragraph()
    p.text = "Implement automated alerts for funds breaching VaR thresholds."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Target SIP investors at the 11-month mark to prevent drop-off."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Diversify away from high HHI (Sector-concentrated) funds."
    p.level = 1

    # --- Slide 12: Thank You ---
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Thank You!"
    subtitle.text = "Questions & Answers\nGitHub: Subhasish-33/bluestock-mutual-fund-analytics"

    # Save presentation
    output_path = os.path.join(os.path.dirname(os.path.dirname(__file__)), 'Bluestock_MF_Presentation.pptx')
    prs.save(output_path)
    print(f"Presentation saved successfully to {output_path}")

if __name__ == "__main__":
    create_presentation()
